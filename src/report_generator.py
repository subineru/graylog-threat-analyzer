"""
Report Generator

Produces a professional blue-themed PowerPoint (.pptx) from aggregated audit data.
Slide deck: Cover → Summary → Distribution → Top Threats → Blocked → Pending → Trend → Footer
"""

import io
from datetime import datetime, timezone

import lxml.etree as etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Cm, Pt

# ---------------------------------------------------------------------------
# Color palette
# ---------------------------------------------------------------------------
NAVY      = RGBColor(0x1B, 0x3A, 0x6B)   # dark navy  — slide headers, cover bg
BLUE      = RGBColor(0x2E, 0x74, 0xB5)   # medium blue — subheadings, borders
ACCENT    = RGBColor(0x44, 0x72, 0xC4)   # PowerPoint blue accent
LIGHT_BG  = RGBColor(0xBD, 0xD7, 0xEE)  # light blue  — table header rows
ALT_ROW   = RGBColor(0xDE, 0xEA, 0xF1)  # pale blue   — alternate table rows
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
GRAY      = RGBColor(0x59, 0x59, 0x59)
COL_RED   = RGBColor(0xC0, 0x00, 0x00)   # block
COL_ORG   = RGBColor(0xED, 0x7D, 0x31)  # monitor / investigate
COL_GRN   = RGBColor(0x37, 0x56, 0x23)  # suppress
COL_BLU   = RGBColor(0x20, 0x6C, 0xD0)  # neutral action count

# Slide dimensions: 16:9  (33.867 cm × 19.05 cm)
W = Cm(33.867)
H = Cm(19.05)

HEADER_H  = Cm(1.9)
MARGIN_L  = Cm(1.5)
MARGIN_R  = Cm(1.5)
CONTENT_W = W - MARGIN_L - MARGIN_R
CONTENT_Y = Cm(2.5)
CONTENT_H = H - CONTENT_Y - Cm(0.8)


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------

def _rgb(cell, color: RGBColor):
    """Set solid fill color on a table cell via XML."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    # remove existing solidFill if any
    for old in tcPr.findall(qn("a:solidFill")):
        tcPr.remove(old)
    solid = etree.SubElement(tcPr, qn("a:solidFill"))
    srgb = etree.SubElement(solid, qn("a:srgbClr"))
    srgb.set("val", f"{color[0]:02X}{color[1]:02X}{color[2]:02X}")


def _para(tf, text: str, size: int, bold=False, color=DARK_TEXT,
          align=PP_ALIGN.LEFT, space_before=0):
    """Add a paragraph to a text frame with consistent styling."""
    p = tf.add_paragraph()
    p.text = text
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def _cell_text(cell, text: str, size: int = 10, bold=False,
               color=DARK_TEXT, align=PP_ALIGN.LEFT):
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text)
    p.alignment = align
    if p.runs:
        run = p.runs[0]
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def _add_slide(prs: Presentation, layout_idx: int = 6):
    """Add a blank slide (layout 6 = blank)."""
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)


def _section_header(slide, title: str):
    """Navy header bar spanning full width with white title text."""
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, W, HEADER_H
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    tf = bar.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = WHITE

    tf.margin_left  = Cm(1.2)
    tf.margin_top   = Cm(0.35)
    tf.margin_bottom = Cm(0)


def _bar_str(value: int, total: int, width: int = 20) -> str:
    """Return a plain-text proportional bar: ████░░░░"""
    if total == 0:
        return "░" * width
    filled = round(value / total * width)
    return "█" * filled + "░" * (width - filled)


def _action_color(action: str) -> RGBColor:
    return {
        "block":       COL_RED,
        "investigate": COL_ORG,
        "monitor":     COL_ORG,
        "suppress":    COL_GRN,
    }.get(action, GRAY)


# ---------------------------------------------------------------------------
# Individual slides
# ---------------------------------------------------------------------------

def _slide_cover(prs: Presentation, stats: dict):
    slide = _add_slide(prs)

    # Navy background
    bg = slide.shapes.add_shape(1, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    # Accent stripe — medium blue horizontal line
    stripe = slide.shapes.add_shape(1, 0, Cm(10.5), W, Cm(0.15))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ACCENT
    stripe.line.fill.background()

    # Title
    tb = slide.shapes.add_textbox(Cm(2), Cm(4), Cm(28), Cm(3.5))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = "Threat Intelligence Report"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = WHITE

    # Subtitle: period
    period = stats.get("period", {})
    sub_text = f"{period.get('start', '')}  ～  {period.get('end', '')}"
    tb2 = slide.shapes.add_textbox(Cm(2), Cm(7.8), Cm(28), Cm(1.2))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = sub_text
    run2 = p2.runs[0]
    run2.font.size = Pt(20)
    run2.font.color.rgb = LIGHT_BG

    # Generated timestamp
    gen_time = datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M UTC")
    tb3 = slide.shapes.add_textbox(Cm(2), Cm(16.5), Cm(28), Cm(0.8))
    tf3 = tb3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = f"Generated: {gen_time}  |  Graylog Threat Analyzer"
    run3 = p3.runs[0]
    run3.font.size = Pt(11)
    run3.font.color.rgb = GRAY


def _slide_summary(prs: Presentation, stats: dict):
    slide = _add_slide(prs)
    _section_header(slide, "Executive Summary")

    total    = stats.get("total_events", 0)
    sup_rt   = stats.get("suppression_rate", 0)
    ac       = stats.get("action_counts", {})
    blocked  = stats.get("edl_active_count", ac.get("block", 0))
    wl_count = stats.get("whitelist_count", 0)

    kpis = [
        ("Total Events",    str(total),           BLUE,    "Events processed in the period"),
        ("EDL 封鎖條目",    str(blocked),          COL_RED, "Active EDL blocked entries"),
        ("Suppression Rate",f"{sup_rt}%",          COL_GRN, "Auto-suppressed without alert"),
        ("白名單規則",      str(wl_count),         COL_BLU, "Known FP whitelist rules"),
    ]

    box_w = Cm(7.2)
    box_h = Cm(5.5)
    gap   = Cm(0.8)
    start_x = Cm(1.2)
    start_y = Cm(3.0)

    for i, (label, value, color, desc) in enumerate(kpis):
        col = i % 2
        row = i // 2
        x = start_x + col * (box_w + gap + Cm(1.4))
        y = start_y + row * (box_h + gap)

        # box border
        box = slide.shapes.add_shape(1, x, y, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = color
        box.line.width = Pt(2)

        # top accent bar
        accent_bar = slide.shapes.add_shape(1, x, y, box_w, Cm(0.35))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = color
        accent_bar.line.fill.background()

        # value
        val_tb = slide.shapes.add_textbox(x + Cm(0.3), y + Cm(0.6), box_w - Cm(0.6), Cm(2.4))
        tf = val_tb.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.color.rgb = color

        # label
        lbl_tb = slide.shapes.add_textbox(x + Cm(0.3), y + Cm(3.1), box_w - Cm(0.6), Cm(0.8))
        tf2 = lbl_tb.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = label
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.runs[0]
        run2.font.size = Pt(13)
        run2.font.bold = True
        run2.font.color.rgb = DARK_TEXT

        # description
        desc_tb = slide.shapes.add_textbox(x + Cm(0.3), y + Cm(4.1), box_w - Cm(0.6), Cm(0.8))
        tf3 = desc_tb.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = desc
        p3.alignment = PP_ALIGN.CENTER
        run3 = p3.runs[0]
        run3.font.size = Pt(9)
        run3.font.color.rgb = GRAY

    # right side: verdict breakdown text
    v = stats.get("verdict_counts", {})
    right_x = start_x + 2 * (box_w + gap + Cm(1.4))
    right_tb = slide.shapes.add_textbox(right_x, start_y, Cm(10), Cm(11.5))
    tf = right_tb.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = "Verdict Breakdown"
    run0 = p0.runs[0]
    run0.font.size = Pt(12)
    run0.font.bold = True
    run0.font.color.rgb = NAVY

    for verdict, count in sorted(v.items(), key=lambda x: -x[1]):
        pct = round(count / total * 100, 1) if total > 0 else 0
        p = tf.add_paragraph()
        p.text = f"  {verdict:<16} {count:>5}  ({pct}%)"
        p.space_before = Pt(6)
        if p.runs:
            p.runs[0].font.size = Pt(11)
            p.runs[0].font.color.rgb = DARK_TEXT


def _slide_distribution(prs: Presentation, stats: dict):
    slide = _add_slide(prs)
    _section_header(slide, "Action Distribution")

    total = stats.get("total_events", 0)
    ac    = stats.get("action_counts", {})

    actions = [
        ("block",       "Confirmed Threat — EDL Block"),
        ("investigate", "Requires Investigation"),
        ("monitor",     "Under Observation"),
        ("suppress",    "Auto-suppressed (Normal / FP)"),
    ]

    # Table: Action | Count | % | Bar
    rows = len(actions) + 1
    cols = 4
    tbl = slide.shapes.add_table(rows, cols,
                                  MARGIN_L, CONTENT_Y,
                                  CONTENT_W, Cm(9)).table

    col_widths = [Cm(5), Cm(3.5), Cm(3), Cm(17)]
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w

    # Header row
    headers = ["Action", "Count", "Percentage", "Distribution"]
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        _rgb(cell, NAVY)
        _cell_text(cell, h, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for ri, (action, desc) in enumerate(actions, start=1):
        count = ac.get(action, 0)
        pct   = round(count / total * 100, 1) if total > 0 else 0
        bar   = _bar_str(count, total, width=30)
        color = ALT_ROW if ri % 2 == 0 else WHITE
        acolor = _action_color(action)

        for ci in range(cols):
            cell = tbl.cell(ri, ci)
            _rgb(cell, color)

        _cell_text(tbl.cell(ri, 0), action.upper(), size=11, bold=True, color=acolor)
        _cell_text(tbl.cell(ri, 1), str(count),     size=11, color=DARK_TEXT, align=PP_ALIGN.CENTER)
        _cell_text(tbl.cell(ri, 2), f"{pct}%",       size=11, color=DARK_TEXT, align=PP_ALIGN.CENTER)
        _cell_text(tbl.cell(ri, 3), bar,             size=9,  color=acolor)

    # Description below table
    desc_tb = slide.shapes.add_textbox(MARGIN_L, Cm(14.5), CONTENT_W, Cm(1.2))
    tf = desc_tb.text_frame
    p = tf.paragraphs[0]
    p.text = f"Total events in period: {total:,}   |   Period: {stats['period']['start']} to {stats['period']['end']}"
    if p.runs:
        p.runs[0].font.size = Pt(10)
        p.runs[0].font.color.rgb = GRAY


def _slide_top_signatures(prs: Presentation, stats: dict):
    slide = _add_slide(prs)
    _section_header(slide, "Top 10 Threat Signatures")

    total = stats.get("total_events", 0)
    top   = stats.get("top_signatures", [])

    rows = min(len(top), 10) + 1
    cols = 4
    tbl = slide.shapes.add_table(rows, cols,
                                  MARGIN_L, CONTENT_Y,
                                  CONTENT_W, Cm(1.1 * rows)).table

    col_widths = [Cm(2.5), Cm(20), Cm(4), Cm(3)]
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w

    for ci, h in enumerate(["Rank", "Signature Name", "Count", "% of Total"]):
        cell = tbl.cell(0, ci)
        _rgb(cell, NAVY)
        _cell_text(cell, h, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for ri, (sig, count) in enumerate(top[:10], start=1):
        pct   = round(count / total * 100, 1) if total > 0 else 0
        color = ALT_ROW if ri % 2 == 0 else WHITE
        for ci in range(cols):
            _rgb(tbl.cell(ri, ci), color)
        _cell_text(tbl.cell(ri, 0), f"#{ri}",     size=10, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
        _cell_text(tbl.cell(ri, 1), sig,           size=10, color=DARK_TEXT)
        _cell_text(tbl.cell(ri, 2), str(count),    size=10, color=DARK_TEXT, align=PP_ALIGN.CENTER)
        _cell_text(tbl.cell(ri, 3), f"{pct}%",     size=10, color=GRAY,      align=PP_ALIGN.CENTER)


def _slide_events_table(prs: Presentation, events: list[dict],
                         title: str, header_color: RGBColor):
    slide = _add_slide(prs)
    _section_header(slide, title)

    if not events:
        tb = slide.shapes.add_textbox(MARGIN_L, CONTENT_Y, CONTENT_W, Cm(3))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = "No events in this category for the selected period."
        if p.runs:
            p.runs[0].font.size = Pt(14)
            p.runs[0].font.color.rgb = GRAY
        return

    show = events[:15]
    rows = len(show) + 1
    cols = 5
    tbl = slide.shapes.add_table(rows, cols,
                                  MARGIN_L, CONTENT_Y,
                                  CONTENT_W, Cm(1.05 * rows)).table

    col_widths = [Cm(4), Cm(5), Cm(5), Cm(8), Cm(10)]
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w

    for ci, h in enumerate(["Timestamp", "Source IP", "Dest IP", "Signature", "Reasoning"]):
        cell = tbl.cell(0, ci)
        _rgb(cell, header_color)
        _cell_text(cell, h, size=10, bold=True, color=WHITE)

    for ri, ev in enumerate(show, start=1):
        color = ALT_ROW if ri % 2 == 0 else WHITE
        for ci in range(cols):
            _rgb(tbl.cell(ri, ci), color)
        _cell_text(tbl.cell(ri, 0), ev.get("timestamp", "")[:16],  size=9, color=DARK_TEXT)
        _cell_text(tbl.cell(ri, 1), ev.get("src_ip", ""),           size=9, color=DARK_TEXT)
        _cell_text(tbl.cell(ri, 2), ev.get("dst_ip", ""),           size=9, color=DARK_TEXT)
        _cell_text(tbl.cell(ri, 3), ev.get("signature", "")[:60],  size=9, color=DARK_TEXT)
        _cell_text(tbl.cell(ri, 4), ev.get("reasoning", "")[:120], size=8, color=GRAY)

    if len(events) > 15:
        note_tb = slide.shapes.add_textbox(MARGIN_L, H - Cm(1.2), CONTENT_W, Cm(0.8))
        tf = note_tb.text_frame
        p = tf.paragraphs[0]
        p.text = f"Showing 15 of {len(events)} events. Full data available via /audit/export."
        if p.runs:
            p.runs[0].font.size = Pt(9)
            p.runs[0].font.color.rgb = GRAY


def _slide_daily_trend(prs: Presentation, stats: dict):
    slide = _add_slide(prs)
    _section_header(slide, "Daily Event Trend")

    daily = stats.get("daily_counts", {})
    if not daily:
        tb = slide.shapes.add_textbox(MARGIN_L, CONTENT_Y, CONTENT_W, Cm(2))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = "No daily data available."
        if p.runs:
            p.runs[0].font.size = Pt(14)
            p.runs[0].font.color.rgb = GRAY
        return

    max_count = max(daily.values()) if daily else 1
    items = sorted(daily.items())

    rows = len(items) + 1
    cols = 3
    tbl = slide.shapes.add_table(rows, cols,
                                  MARGIN_L, CONTENT_Y,
                                  CONTENT_W, Cm(0.85 * rows)).table

    col_widths = [Cm(5), Cm(4), Cm(24)]
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w

    for ci, h in enumerate(["Date", "Events", "Trend"]):
        cell = tbl.cell(0, ci)
        _rgb(cell, NAVY)
        _cell_text(cell, h, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for ri, (day, count) in enumerate(items, start=1):
        color = ALT_ROW if ri % 2 == 0 else WHITE
        bar = _bar_str(count, max_count, width=35)
        for ci in range(cols):
            _rgb(tbl.cell(ri, ci), color)
        _cell_text(tbl.cell(ri, 0), day,        size=10, color=DARK_TEXT, align=PP_ALIGN.CENTER)
        _cell_text(tbl.cell(ri, 1), str(count), size=10, color=DARK_TEXT, align=PP_ALIGN.CENTER)
        _cell_text(tbl.cell(ri, 2), bar,        size=9,  color=ACCENT)


def _slide_footer(prs: Presentation, stats: dict):
    slide = _add_slide(prs)

    bg = slide.shapes.add_shape(1, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    stripe = slide.shapes.add_shape(1, 0, Cm(9.2), W, Cm(0.12))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ACCENT
    stripe.line.fill.background()

    tb = slide.shapes.add_textbox(Cm(2), Cm(6), Cm(28), Cm(3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "End of Report"
    p.alignment = PP_ALIGN.CENTER
    if p.runs:
        p.runs[0].font.size = Pt(28)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = WHITE

    tb2 = slide.shapes.add_textbox(Cm(2), Cm(10), Cm(28), Cm(1.5))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    period = stats.get("period", {})
    p2.text = f"Graylog Threat Analyzer  |  {period.get('start', '')} – {period.get('end', '')}"
    p2.alignment = PP_ALIGN.CENTER
    if p2.runs:
        p2.runs[0].font.size = Pt(13)
        p2.runs[0].font.color.rgb = LIGHT_BG


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------

def generate_pptx(stats: dict) -> bytes:
    """
    Build a professional blue-themed .pptx from aggregated audit stats.
    Returns raw bytes of the .pptx file.
    """
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    _slide_cover(prs, stats)
    _slide_summary(prs, stats)
    _slide_distribution(prs, stats)
    _slide_top_signatures(prs, stats)
    _slide_events_table(prs, stats.get("block_events", []),
                        "封鎖建議事件 (Block-Recommended by Triage)", COL_RED)
    _slide_events_table(prs, stats.get("pending_events", []),
                        "Pending Review — Monitor / Investigate", COL_ORG)
    _slide_daily_trend(prs, stats)
    _slide_footer(prs, stats)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
